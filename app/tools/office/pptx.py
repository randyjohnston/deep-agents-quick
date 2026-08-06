"""Typed PPTX generation with bounded slide archetypes and text extraction."""

from __future__ import annotations

from pathlib import Path
from typing import Annotated, Literal

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, ConfigDict, Field

from app.tools.office.paths import (
    load_office_file,
    load_office_template,
    resolve_image_path,
    resolve_read_path,
    resolve_template_path,
    resolve_write_path,
)
from app.tools.office.pptx_grid import Grid, WIDE_HEIGHT, WIDE_WIDTH
from app.tools.office.theme import ResolvedTheme, Theme, resolve_theme

_POTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
)
_PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)


class _PptxModel(BaseModel):
    """Forbid unbounded layout fields across the entire PPTX schema."""

    model_config = ConfigDict(extra="forbid", str_strip_whitespace=True)


class Slide(_PptxModel):
    """Backward-compatible slide with a title and concise bullet points."""

    title: str = Field(max_length=100, description="Slide title")
    bullets: list[str] = Field(default_factory=list, max_length=8)


class CoverSlide(_PptxModel):
    """Full-bleed opening slide with optional photography."""

    kind: Literal["cover"]
    title: str = Field(max_length=70)
    subtitle: str | None = Field(default=None, max_length=160)
    kicker: str | None = Field(default=None, max_length=50)
    meta: str | None = Field(default=None, max_length=80)
    image: str | None = Field(default=None, max_length=255)
    notes: str | None = Field(default=None, max_length=4000)


class Stat(_PptxModel):
    """One bounded metric shown on a stat-card slide."""

    value: str = Field(max_length=18)
    label: str = Field(max_length=45)
    caption: str | None = Field(default=None, max_length=120)


class StatsSlide(_PptxModel):
    """A headline followed by two to four metric cards."""

    kind: Literal["stats"]
    headline: str = Field(max_length=75)
    kicker: str | None = Field(default=None, max_length=50)
    deck: str | None = Field(default=None, max_length=160)
    stats: list[Stat] = Field(min_length=2, max_length=4)
    source: str | None = Field(default=None, max_length=180)
    notes: str | None = Field(default=None, max_length=4000)


class Card(_PptxModel):
    """One image card with tightly bounded display copy."""

    title: str = Field(max_length=48)
    metric: str | None = Field(default=None, max_length=45)
    body: str = Field(max_length=120)
    image: str | None = Field(default=None, max_length=255)


class CardsSlide(_PptxModel):
    """A two- to four-column image-card composition."""

    kind: Literal["cards"]
    headline: str = Field(max_length=75)
    kicker: str | None = Field(default=None, max_length=50)
    deck: str | None = Field(default=None, max_length=160)
    cards: list[Card] = Field(min_length=2, max_length=4)
    source: str | None = Field(default=None, max_length=180)
    notes: str | None = Field(default=None, max_length=4000)


class StatementSlide(_PptxModel):
    """Oversized closing statement with optional attribution."""

    kind: Literal["statement"]
    statement: str = Field(max_length=140)
    kicker: str | None = Field(default=None, max_length=50)
    attribution: str | None = Field(default=None, max_length=90)
    notes: str | None = Field(default=None, max_length=4000)


ArchetypeSlide = Annotated[
    CoverSlide | StatsSlide | CardsSlide | StatementSlide,
    Field(discriminator="kind"),
]


class Deck(_PptxModel):
    """A PowerPoint deck supporting legacy bullets and designed archetypes."""

    title: str | None = Field(default=None, max_length=90)
    subtitle: str | None = Field(default=None, max_length=160)
    slides: list[Slide | ArchetypeSlide] = Field(default_factory=list, max_length=100)


class _Design:
    """Resolved PPTX design tokens with conservative defaults."""

    def __init__(self, theme: ResolvedTheme | None):
        self.ink = (theme.ink_color if theme else None) or "111111"
        self.muted = (theme.muted_color if theme else None) or "555A62"
        self.accent = (theme.accent_color if theme else None) or "336699"
        self.surface = (theme.surface_color if theme else None) or "F2F2F2"
        self.on_accent = (theme.on_accent_color if theme else None) or "FFFFFF"
        self.rule = (theme.rule_color if theme else None) or "B8BCC4"
        self.heading_font = (theme.heading_font if theme else None) or "Aptos Display"
        self.body_font = (theme.body_font if theme else None) or "Aptos"
        self.cover_size = (theme.cover_title_size if theme else None) or 61
        self.stat_size = (theme.stat_size if theme else None) or 35
        self.headline_size = (theme.headline_size if theme else None) or 21
        self.deck_size = (theme.deck_size if theme else None) or 15
        self.body_size = (theme.body_size if theme else None) or 13
        self.small_size = (theme.small_size if theme else None) or 9
        self.fine_size = (theme.fine_size if theme else None) or 7


def write_pptx(
    filename: str,
    deck: Deck,
    theme: Theme | None = None,
    template: str | None = None,
) -> str:
    """Create a .pptx from legacy bullets or bounded designed archetypes."""
    if not deck.title and not deck.subtitle and not deck.slides:
        raise ValueError("Provide a title, subtitle, or at least one slide.")
    path = resolve_write_path(filename, ".pptx")
    path.parent.mkdir(parents=True, exist_ok=True)
    template_path = resolve_template_path(template, (".potx", ".pptx")) if template else None
    presentation = (
        load_office_template(template_path, Presentation, _POTX_CONTENT_TYPE, _PPTX_CONTENT_TYPE)
        if template_path
        else Presentation()
    )
    if template_path is None:
        presentation.slide_width = WIDE_WIDTH
        presentation.slide_height = WIDE_HEIGHT
    branding = resolve_theme(theme)
    design = _Design(branding)
    if deck.title or deck.subtitle:
        _render_legacy_title(presentation, deck.title, deck.subtitle, branding)
    for index, item in enumerate(deck.slides, start=1):
        if isinstance(item, Slide):
            _render_bullets(presentation, item, branding)
        elif isinstance(item, CoverSlide):
            _render_cover(presentation, item, design, branding)
        elif isinstance(item, StatsSlide):
            _render_stats(presentation, item, design, branding, index)
        elif isinstance(item, CardsSlide):
            _render_cards(presentation, item, design, branding, index)
        else:
            _render_statement(presentation, item, design, branding)
    presentation.save(path)
    return f"Wrote {path} — {len(presentation.slides)} slides"


def _render_legacy_title(presentation, title, subtitle, theme) -> None:
    layout, body_idx = _content_layout(presentation, "Title Slide", 0)
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = title or ""
    _placeholder(slide, body_idx).text = subtitle or ""
    _style_slide(slide, theme)
    _add_logo(slide, presentation, theme)


def _render_bullets(presentation, item: Slide, theme) -> None:
    layout, body_idx = _content_layout(presentation, "Title and Content", 1)
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = item.title
    frame = _placeholder(slide, body_idx).text_frame
    frame.clear()
    for index, bullet in enumerate(item.bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
    _style_slide(slide, theme)
    _add_logo(slide, presentation, theme)


def _render_cover(presentation, item, design, theme) -> None:
    slide = _blank_slide(presentation)
    grid = Grid(presentation.slide_width, presentation.slide_height)
    width, height = grid.width, grid.height
    if item.image:
        _add_fill_image(slide, resolve_image_path(item.image), 0, 0, width, height)
    else:
        _rect(slide, 0, 0, width, height, design.surface)
    panel_width = grid.x(.58)
    _rect(slide, 0, 0, panel_width, height, design.ink)
    left, content_width = grid.x(.055), int(panel_width * .82)
    _text(slide, item.kicker, left, grid.y(.10), content_width, Inches(.3),
          design.small_size, design.accent, design.body_font, bold=True)
    _text(slide, item.title, left, grid.y(.23), content_width, Inches(2.2),
          design.cover_size, design.on_accent, design.heading_font, bold=True)
    _text(slide, item.subtitle, left, grid.y(.61), content_width, Inches(.9),
          design.deck_size, design.on_accent, design.body_font)
    _text(slide, item.meta, left, grid.y(.84), content_width, Inches(.35),
          design.small_size, design.rule, design.body_font)
    _add_logo(slide, presentation, theme)
    _notes(slide, item.notes)


def _render_stats(presentation, item, design, theme, page) -> None:
    slide = _blank_slide(presentation)
    _chrome(slide, presentation, item.kicker, item.headline, item.deck, item.source, page, design)
    grid = Grid(presentation.slide_width, presentation.slide_height)
    top, card_height = grid.y(.27), grid.y(.50)
    for stat, (left, card_width) in zip(item.stats, grid.columns(len(item.stats)), strict=True):
        _rect(slide, left, top, card_width, card_height, design.surface)
        _rect(slide, left, top, Inches(.07), card_height, design.accent)
        _text(slide, stat.value, left + Inches(.25), top + Inches(.30), card_width - Inches(.5),
              Inches(.65), design.stat_size, design.accent, design.heading_font, bold=True)
        _text(slide, stat.label, left + Inches(.25), top + Inches(1.15), card_width - Inches(.5),
              Inches(.55), design.deck_size, design.ink, design.heading_font, bold=True)
        _text(slide, stat.caption, left + Inches(.25), top + Inches(1.85), card_width - Inches(.5),
              Inches(1.2), design.body_size, design.muted, design.body_font)
    _add_logo(slide, presentation, theme)
    _notes(slide, item.notes)


def _render_cards(presentation, item, design, theme, page) -> None:
    slide = _blank_slide(presentation)
    _chrome(slide, presentation, item.kicker, item.headline, item.deck, item.source, page, design)
    grid = Grid(presentation.slide_width, presentation.slide_height)
    top, image_height = grid.y(.25), grid.y(.29)
    for card, (left, card_width) in zip(item.cards, grid.columns(len(item.cards)), strict=True):
        _rect(slide, left, top, card_width, grid.y(.53), design.surface)
        if card.image:
            _add_fill_image(
                slide, resolve_image_path(card.image), left, top, card_width, image_height
            )
        else:
            _rect(slide, left, top, card_width, image_height, design.surface)
            _rect(
                slide, left, top + image_height - Inches(.07), card_width,
                Inches(.07), design.accent,
            )
        pad = Inches(.18)
        _text(
            slide, card.title, left + pad, top + image_height + Inches(.15),
            card_width - 2 * pad, Inches(.45), design.deck_size, design.ink,
            design.heading_font, bold=True,
        )
        _text(
            slide, card.metric, left + pad, top + image_height + Inches(.63),
            card_width - 2 * pad, Inches(.34), design.small_size, design.accent,
            design.body_font, bold=True,
        )
        _text(slide, card.body, left + pad, top + image_height + Inches(1.02), card_width - 2 * pad,
              Inches(.82), design.body_size, design.muted, design.body_font)
    _add_logo(slide, presentation, theme)
    _notes(slide, item.notes)


def _render_statement(presentation, item, design, theme) -> None:
    slide = _blank_slide(presentation)
    grid = Grid(presentation.slide_width, presentation.slide_height)
    width, height = grid.width, grid.height
    _rect(slide, 0, 0, width, height, design.ink)
    _rect(slide, grid.x(.055), grid.y(.18), Inches(.08), grid.y(.52), design.accent)
    _text(slide, item.kicker, grid.x(.085), grid.y(.17), grid.x(.75), Inches(.35),
          design.small_size, design.accent, design.body_font, bold=True)
    _text(slide, item.statement, grid.x(.085), grid.y(.29), grid.x(.76), Inches(2.5),
          42, design.on_accent, design.heading_font, bold=True)
    _text(slide, item.attribution, grid.x(.085), grid.y(.76), grid.x(.70), Inches(.4),
          design.body_size, design.rule, design.body_font)
    _add_logo(slide, presentation, theme)
    _notes(slide, item.notes)


def _chrome(slide, presentation, kicker, headline, deck, source, page, design) -> None:
    grid = Grid(presentation.slide_width, presentation.slide_height)
    width, height, margin = grid.width, grid.height, grid.margin
    _text(slide, kicker, margin, grid.y(.035), grid.x(.75), Inches(.25),
          design.small_size, design.accent, design.body_font, bold=True)
    _text(slide, headline, margin, grid.y(.075), grid.x(.80), Inches(.55),
          design.headline_size, design.ink, design.heading_font, bold=True)
    _text(slide, deck, margin, grid.y(.15), grid.x(.82), Inches(.4),
          design.deck_size, design.muted, design.body_font)
    _rect(slide, margin, grid.y(.205), width - 2 * margin, max(1, Inches(.01)), design.rule)
    _rect(slide, margin, grid.y(.908), width - 2 * margin, max(1, Inches(.01)), design.rule)
    _text(slide, source, margin, grid.y(.92), grid.x(.75), Inches(.22),
          design.fine_size, design.muted, design.body_font)
    _text(slide, str(page), grid.x(.92), grid.y(.92), grid.x(.04), Inches(.22),
          design.fine_size, design.muted, design.body_font, align=PP_ALIGN.RIGHT)


def _blank_slide(presentation):
    layout = next(
        (
            candidate
            for candidate in presentation.slide_layouts
            if candidate.name.casefold() == "blank"
        ),
        min(presentation.slide_layouts, key=lambda item: len(item.placeholders)),
    )
    slide = presentation.slides.add_slide(layout)
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            shape._element.getparent().remove(shape._element)
    return slide


def _add_fill_image(slide, path: Path, left, top, width, height) -> None:
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    frame_ratio = width / height
    picture = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if image_ratio > frame_ratio:
        crop = (1 - frame_ratio / image_ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif image_ratio < frame_ratio:
        crop = (1 - image_ratio / frame_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop


def _rect(slide, left, top, width, height, color) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()


def _text(slide, value, left, top, width, height, size, color, font, bold=False, align=None):
    if not value:
        return
    frame = slide.shapes.add_textbox(left, top, width, height).text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _notes(slide, notes: str | None) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _content_layout(presentation, preferred_name: str, fallback_index: int):
    layout = next(
        (candidate for candidate in presentation.slide_layouts
         if candidate.name.casefold() == preferred_name.casefold()),
        None,
    )
    if layout is None and len(presentation.slide_layouts) <= fallback_index:
        raise ValueError(f"Presentation template has no usable {preferred_name!r} layout")
    layout = layout or presentation.slide_layouts[fallback_index]
    title = next(
        (
            placeholder for placeholder in layout.placeholders
            if placeholder.placeholder_format.type
            in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
        ),
        None,
    )
    body = next(
        (placeholder for placeholder in layout.placeholders
         if title and placeholder.shape_id != title.shape_id
         and getattr(placeholder, "has_text_frame", False)
         and placeholder.placeholder_format.type not in {
             PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.SLIDE_NUMBER}),
        None,
    )
    if title is None or body is None:
        raise ValueError(
            f"Presentation template layout {layout.name!r} must contain title and body placeholders"
        )
    return layout, body.placeholder_format.idx


def _placeholder(slide, idx: int):
    placeholder = next(
        (candidate for candidate in slide.placeholders if candidate.placeholder_format.idx == idx),
        None,
    )
    if placeholder is None:
        raise ValueError(f"Presentation template did not clone required placeholder {idx}")
    return placeholder


def _style_slide(slide, theme: ResolvedTheme | None) -> None:
    if not theme:
        return
    title = slide.shapes.title
    if title:
        if theme.header_background:
            title.fill.solid()
            title.fill.fore_color.rgb = RGBColor.from_string(theme.header_background)
        _style_text(title.text_frame, theme.heading_font, theme.header_foreground, bold=True)
    for shape in slide.placeholders:
        if (title and shape.shape_id == title.shape_id) or not getattr(
            shape, "has_text_frame", False
        ):
            continue
        _style_text(shape.text_frame, theme.body_font, theme.accent_color)


def _style_text(frame, font_name, color, bold=False) -> None:
    for paragraph in frame.paragraphs:
        for run in paragraph.runs:
            if font_name:
                run.font.name = font_name
            if color:
                run.font.color.rgb = RGBColor.from_string(color)
            if bold:
                run.font.bold = True


def _add_logo(slide, presentation, theme: ResolvedTheme | None) -> None:
    if theme and theme.logo:
        width = Inches(1)
        slide.shapes.add_picture(
            str(theme.logo), presentation.slide_width - width - Inches(.25), Inches(.1), width=width
        )


def read_pptx(path: str, max_slides: int = 100) -> str:
    """Extract visible text from a permitted .pptx file."""
    if max_slides < 1:
        raise ValueError("max_slides must be at least 1.")
    resolved = resolve_read_path(path, (".pptx",))
    presentation = load_office_file(resolved, Presentation)
    blocks = [f"# {resolved}", ""]
    for index, slide in enumerate(presentation.slides, start=1):
        if index > max_slides:
            break
        blocks.append(f"## Slide {index}")
        text = [value for shape in slide.shapes for value in _shape_text(shape)]
        blocks.extend(text or ["(empty)"])
        blocks.append("")
    if len(presentation.slides) > max_slides:
        blocks.append(f"... truncated at {max_slides} slides; raise max_slides for more")
    return "\n".join(blocks).rstrip()


def _shape_text(shape) -> list[str]:
    if getattr(shape, "has_table", False):
        return [
            " | ".join(cell.text for cell in row.cells)
            for row in shape.table.rows
            if any(cell.text.strip() for cell in row.cells)
        ]
    if hasattr(shape, "shapes"):
        return [value for child in shape.shapes for value in _shape_text(child)]
    text = getattr(shape, "text", "")
    return [text] if text.strip() else []
