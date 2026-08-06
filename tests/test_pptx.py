from __future__ import annotations

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.tools.office.pptx import (
    Card, CardsSlide, CoverSlide, Deck, Slide, Stat, StatementSlide, StatsSlide,
    read_pptx, write_pptx,
)
from app.tools.office.theme import Theme


def test_pptx_round_trip_with_title_and_bullets(_isolate_dirs):
    result = write_pptx(
        "deck",
        Deck(
            title="Quarterly Review",
            subtitle="Q3",
            slides=[Slide(title="Growth", bullets=["18% YoY"])],
        ),
    )
    path = _isolate_dirs / "out" / "deck.pptx"
    assert path.is_file(), result
    parsed = Presentation(path)
    assert len(parsed.slides) == 2
    assert parsed.slides[1].shapes.title.text == "Growth"
    assert "18% YoY" in read_pptx("deck.pptx")


def test_pptx_rejects_empty_content(_isolate_dirs):
    with pytest.raises(ValueError, match="Provide"):
        write_pptx("empty", Deck())


def test_pptx_rejects_macro_extension(_isolate_dirs):
    deck = Deck(title="Safe")
    with pytest.raises(ValueError, match="Macro-enabled"):
        write_pptx("unsafe.pptm", deck)


def test_pptx_read_includes_table_text(_isolate_dirs):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Q3 Numbers"
    table = slide.shapes.add_table(1, 2, 0, 0, Inches(4), Inches(1)).table
    table.cell(0, 0).text = "EMEA"
    table.cell(0, 1).text = "4200000"
    path = _isolate_dirs / "in" / "table.pptx"
    presentation.save(path)

    output = read_pptx("table.pptx")
    assert "Q3 Numbers" in output
    assert "EMEA | 4200000" in output


def test_pptx_read_recurses_into_grouped_shapes(_isolate_dirs):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    textbox = group.shapes.add_textbox(0, 0, Inches(4), Inches(1))
    textbox.text = "Nested takeaway"
    path = _isolate_dirs / "in" / "group.pptx"
    presentation.save(path)

    assert "Nested takeaway" in read_pptx("group.pptx")


def test_designed_archetypes_render_wide_without_images(_isolate_dirs):
    write_pptx(
        "designed",
        Deck(slides=[
            CoverSlide(kind="cover", kicker="MOBILITY", title="Built for what is next"),
            StatsSlide(kind="stats", headline="Momentum is measurable", stats=[
                Stat(value="12", label="Launches", caption="Across key segments"),
                Stat(value="3", label="Powertrains", caption="Choice without compromise"),
            ]),
            CardsSlide(kind="cards", headline="A portfolio for every road", cards=[
                Card(title="Sedan", metric="Efficient", body="Confident daily mobility."),
                Card(title="SUV", metric="Capable", body="Space for every adventure."),
            ]),
            StatementSlide(kind="statement", statement="Let us go places together."),
        ]),
        theme=Theme(accent_color="EB0A1E", heading_font="Arial", body_font="Arial"),
    )
    presentation = Presentation(_isolate_dirs / "out" / "designed.pptx")

    assert presentation.slide_width == Inches(13.333333)
    assert presentation.slide_height == Inches(7.5)
    assert len(presentation.slides) == 4
    assert not any(
        shape.shape_type == 13 for slide in presentation.slides for shape in slide.shapes
    )
    assert not any(
        shape.is_placeholder for slide in presentation.slides for shape in slide.shapes
    )
    assert "Built for what is next" in read_pptx("designed.pptx")

    sedan = next(shape for shape in presentation.slides[2].shapes if shape.text == "Sedan")
    assert sedan.top < Inches(3)


def test_image_cards_crop_to_fill_and_notes_round_trip(_isolate_dirs):
    image = _isolate_dirs / "in" / "wide.png"
    Image.new("RGB", (800, 200), "red").save(image)
    write_pptx("photos", Deck(slides=[CardsSlide(
        kind="cards", headline="Designed details", notes="Discuss the launch cadence.",
        cards=[
            Card(title="A", body="First card", image="wide.png"),
            Card(title="B", body="Second card", image="wide.png"),
        ],
    )]))
    slide = Presentation(_isolate_dirs / "out" / "photos.pptx").slides[0]
    pictures = [shape for shape in slide.shapes if shape.shape_type == 13]

    assert len(pictures) == 2
    assert all(picture.crop_left > 0 and picture.crop_right > 0 for picture in pictures)
    assert "Discuss the launch cadence." in slide.notes_slide.notes_text_frame.text


def test_archetype_text_budgets_are_schema_enforced():
    with pytest.raises(ValueError, match="at most 70 characters"):
        CoverSlide(kind="cover", title="x" * 71)
    with pytest.raises(ValueError, match="at most 120 characters"):
        Card(title="Card", body="x" * 121)
    with pytest.raises(ValueError, match="Extra inputs"):
        CoverSlide(kind="cover", title="Bounded", x=10)

    parsed = Deck.model_validate({"slides": [{"kind": "cover", "title": "Typed"}]})
    assert isinstance(parsed.slides[0], CoverSlide)


def test_designed_page_numbers_follow_rendered_slide_positions(_isolate_dirs):
    write_pptx(
        "pages",
        Deck(
            title="Legacy title",
            slides=[
                CoverSlide(kind="cover", title="Cover"),
                StatsSlide(
                    kind="stats", headline="Stats",
                    stats=[Stat(value="A", label="One"), Stat(value="B", label="Two")],
                ),
                CardsSlide(
                    kind="cards", headline="Cards",
                    cards=[Card(title="A", body="One"), Card(title="B", body="Two")],
                ),
                StatementSlide(kind="statement", statement="Close"),
            ],
        ),
    )
    presentation = Presentation(_isolate_dirs / "out" / "pages.pptx")

    assert [_page_number(slide, presentation.slide_width) for slide in presentation.slides] == [
        None, None, "3", "4", "5"
    ]


def _page_number(slide, width):
    values = [shape.text for shape in slide.shapes if shape.left >= int(width * .9)]
    return values[-1] if values else None
