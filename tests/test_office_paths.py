from __future__ import annotations

from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from docx import Document as DocxDocument
from docx.shared import RGBColor
from openpyxl import Workbook, load_workbook
from openpyxl.styles import NamedStyle
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.tools.office import paths
from app.tools.office.docx import Document, read_docx, write_docx
from app.tools.office.pptx import Deck, read_pptx, write_pptx
from app.tools.office.xlsx import Sheet, read_xlsx, write_xlsx
from app.tools.office.theme import Theme, resolve_theme


def _write_xlsx(name, **kwargs):
    return write_xlsx(name, [Sheet(name="S", columns=["C"], rows=[["x"]])], **kwargs)


def _write_docx(name, **kwargs):
    return write_docx(name, Document(title="T", sections=[]), **kwargs)


def _write_pptx(name, **kwargs):
    return write_pptx(name, Deck(title="T"), **kwargs)


WRITER_CASES = [
    pytest.param(_write_xlsx, ".xlsx", id="xlsx"),
    pytest.param(_write_docx, ".docx", id="docx"),
    pytest.param(_write_pptx, ".pptx", id="pptx"),
]
READER_CASES = [
    pytest.param(read_xlsx, ".xlsx", id="xlsx"),
    pytest.param(read_docx, ".docx", id="docx"),
    pytest.param(read_pptx, ".pptx", id="pptx"),
]


@pytest.mark.parametrize("writer,extension", WRITER_CASES)
def test_writers_add_extension_and_reject_escape(_isolate_dirs, writer, extension):
    writer("safe")
    assert (_isolate_dirs / "out" / f"safe{extension}").is_file()
    with pytest.raises(ValueError, match="directories"):
        writer(f"../escape{extension}")


@pytest.mark.parametrize("reader,extension", READER_CASES)
def test_readers_refuse_paths_outside_roots(_isolate_dirs, reader, extension):
    outside = _isolate_dirs / f"outside{extension}"
    outside.write_bytes(b"not opened")
    with pytest.raises(FileNotFoundError):
        reader(str(outside))


def test_read_rejects_expanded_archive_over_limit(_isolate_dirs, monkeypatch):
    archive = _isolate_dirs / "in" / "bomb.docx"
    with ZipFile(archive, "w", ZIP_DEFLATED) as output:
        output.writestr("word/document.xml", b"x" * 10_000)
    monkeypatch.setattr(paths, "MAX_EXPANDED_BYTES", 1_000)
    with pytest.raises(ValueError, match="expanded archive"):
        paths.resolve_read_path("bomb.docx", (".docx",))


def test_read_rejects_compressed_archive_over_limit(_isolate_dirs, monkeypatch):
    archive = _isolate_dirs / "in" / "large.pptx"
    with ZipFile(archive, "w") as output:
        output.writestr("ppt/presentation.xml", b"content")
    monkeypatch.setattr(paths, "MAX_ARCHIVE_BYTES", 1)
    with pytest.raises(ValueError, match="archive is"):
        paths.resolve_read_path("large.pptx", (".pptx",))


def test_read_rejects_invalid_zip(_isolate_dirs):
    invalid = _isolate_dirs / "in" / "invalid.xlsx"
    invalid.write_text("not a zip")
    with pytest.raises(ValueError, match="Not a valid"):
        paths.resolve_read_path("invalid.xlsx", (".xlsx",))


def test_reader_normalizes_corrupt_member_error(_isolate_dirs):
    corrupt = _isolate_dirs / "in" / "corrupt.docx"
    DocxDocument().save(corrupt)
    with ZipFile(corrupt) as archive:
        member = archive.getinfo("word/document.xml")
        payload_offset = (
            member.header_offset + 30 + len(member.filename.encode()) + len(member.extra)
        )
    data = bytearray(corrupt.read_bytes())
    data[payload_offset] ^= 0xFF
    corrupt.write_bytes(data)

    with pytest.raises(ValueError, match="Not a valid Office Open XML"):
        read_docx("corrupt.docx")


@pytest.mark.parametrize("reader,extension", READER_CASES)
def test_readers_normalize_missing_office_parts(_isolate_dirs, reader, extension):
    plain_zip = _isolate_dirs / "in" / f"plain{extension}"
    with ZipFile(plain_zip, "w") as archive:
        archive.writestr("ordinary.txt", "not an Office package")

    with pytest.raises(ValueError, match="Not a valid Office Open XML"):
        reader(plain_zip.name)


@pytest.mark.parametrize(
    "writer,reader,extension,member",
    [
        pytest.param(
            _write_xlsx, read_xlsx, ".xlsx", "xl/workbook.xml", id="xlsx"
        ),
        pytest.param(
            _write_docx, read_docx, ".docx", "word/document.xml", id="docx"
        ),
        pytest.param(
            _write_pptx,
            read_pptx,
            ".pptx",
            "ppt/presentation.xml",
            id="pptx",
        ),
    ],
)
def test_readers_normalize_malformed_xml(_isolate_dirs, writer, reader, extension, member):
    writer("malformed")
    path = _isolate_dirs / "out" / f"malformed{extension}"
    _replace_member(path, member, b"<not valid XML")

    with pytest.raises(ValueError, match="Not a valid Office Open XML"):
        reader(path.name)


def _replace_member(path, target: str, replacement: bytes) -> None:
    buffer = BytesIO()
    with ZipFile(path) as source, ZipFile(buffer, "w") as destination:
        for member in source.infolist():
            destination.writestr(
                member,
                replacement if member.filename == target else source.read(member),
            )
    path.write_bytes(buffer.getvalue())


def test_named_theme_resolves_with_inline_overrides(_isolate_dirs):
    (_isolate_dirs / "themes" / "acme.json").write_text(
        '{"accent_color":"112233","body_font":"Aptos"}'
    )
    resolved = resolve_theme(Theme(name="acme", accent_color="AABBCC"))

    assert resolved.accent_color == "AABBCC"
    assert resolved.body_font == "Aptos"


def test_named_toml_theme_resolves(_isolate_dirs):
    (_isolate_dirs / "themes" / "acme.toml").write_text('heading_font = "Arial"')
    assert resolve_theme(Theme(name="acme")).heading_font == "Arial"


def test_named_theme_is_confined_sized_and_unambiguous(_isolate_dirs, monkeypatch):
    outside = _isolate_dirs / "outside.json"
    outside.write_text('{"accent_color":"112233"}')
    (_isolate_dirs / "themes" / "escape.json").symlink_to(outside)
    with pytest.raises(FileNotFoundError, match="No named Office theme"):
        resolve_theme(Theme(name="escape"))

    json_theme = _isolate_dirs / "themes" / "acme.json"
    json_theme.write_text('{"accent_color":"112233"}')
    monkeypatch.setattr("app.tools.office.theme.MAX_THEME_BYTES", 1)
    with pytest.raises(ValueError, match="exceeds"):
        resolve_theme(Theme(name="acme"))

    monkeypatch.setattr("app.tools.office.theme.MAX_THEME_BYTES", 64 * 1024)
    (_isolate_dirs / "themes" / "acme.toml").write_text('accent_color = "112233"')
    with pytest.raises(ValueError, match="ambiguous"):
        resolve_theme(Theme(name="acme"))


def test_theme_rejects_unbounded_fields_and_invalid_values():
    with pytest.raises(ValueError, match="Extra inputs"):
        Theme.model_validate({"raw_xml": "<xml/>"})
    with pytest.raises(ValueError, match="six hexadecimal"):
        Theme(accent_color="#123456")
    with pytest.raises(ValueError, match="theme name"):
        Theme(name="../acme")


def test_logo_is_confined_and_pixel_bounded(_isolate_dirs, monkeypatch):
    logo = _isolate_dirs / "in" / "logo.png"
    Image.new("RGB", (2, 2), "red").save(logo)
    assert resolve_theme(Theme(logo="logo.png")).logo == logo

    monkeypatch.setattr(paths, "MAX_IMAGE_PIXELS", 3)
    with pytest.raises(ValueError, match="and 3 pixels"):
        resolve_theme(Theme(logo="logo.png"))
    with pytest.raises(FileNotFoundError, match="permitted Office asset"):
        resolve_theme(Theme(logo=str(_isolate_dirs / "outside.png")))


@pytest.mark.parametrize("writer,extension", WRITER_CASES)
def test_named_theme_and_logo_apply_across_formats(_isolate_dirs, writer, extension):
    logo = _isolate_dirs / "in" / "logo.png"
    Image.new("RGB", (20, 10), "blue").save(logo)
    (_isolate_dirs / "themes" / "acme.json").write_text(
        """{
            "accent_color": "336699",
            "header_background": "112233",
            "header_foreground": "FFFFFF",
            "heading_font": "Arial",
            "body_font": "Courier New",
            "logo": "logo.png"
        }"""
    )
    writer("branded", theme=Theme(name="acme"))
    output = _isolate_dirs / "out" / f"branded{extension}"

    if extension == ".xlsx":
        sheet = load_workbook(output).active
        assert sheet["A1"].fill.fgColor.rgb == "FF112233"
        assert sheet["A1"].font.name == "Arial"
        assert sheet["A2"].font.name == "Courier New"
        assert len(sheet._images) == 1
    elif extension == ".docx":
        document = DocxDocument(output)
        assert document.styles["Normal"].font.name == "Courier New"
        assert document.styles["Title"].font.name == "Arial"
        assert document.styles["Title"].font.color.rgb == RGBColor(255, 255, 255)
        assert document.styles["Title"].element.xpath("./w:pPr/w:shd/@w:fill") == ["112233"]
        assert len(document.inline_shapes) == 1
    else:
        presentation = Presentation(output)
        title = presentation.slides[0].shapes.title
        assert title.fill.fore_color.rgb == RGBColor(17, 34, 51)
        assert title.text_frame.paragraphs[0].runs[0].font.name == "Arial"
        assert len(presentation.slides[0].shapes) >= 3


def test_xlsx_template_preserves_named_styles(_isolate_dirs):
    template = Workbook()
    template.add_named_style(NamedStyle(name="AcmeBody"))
    template.template = True
    template.save(_isolate_dirs / "in" / "acme.xltx")
    _write_xlsx("from-template", template="acme.xltx")

    output = load_workbook(_isolate_dirs / "out" / "from-template.xlsx")
    assert "AcmeBody" in output.named_styles
    assert output.template is False


def test_docx_template_preserves_styles(_isolate_dirs):
    template = DocxDocument()
    template.styles["Normal"].font.name = "Courier New"
    template_path = _isolate_dirs / "in" / "acme.dotx"
    template.save(template_path)
    _replace_bytes_in_member(
        template_path,
        "[Content_Types].xml",
        b"wordprocessingml.document.main+xml",
        b"wordprocessingml.template.main+xml",
    )
    _write_docx("from-template", template="acme.dotx")

    output_path = _isolate_dirs / "out" / "from-template.docx"
    output = DocxDocument(output_path)
    assert output.styles["Normal"].font.name == "Courier New"
    with ZipFile(output_path) as package:
        assert b"wordprocessingml.document.main+xml" in package.read("[Content_Types].xml")


def test_pptx_template_preserves_page_size(_isolate_dirs):
    template = Presentation()
    template.slide_width = Inches(12)
    template_path = _isolate_dirs / "in" / "acme.potx"
    template.save(template_path)
    _replace_bytes_in_member(
        template_path,
        "[Content_Types].xml",
        b"presentationml.presentation.main+xml",
        b"presentationml.template.main+xml",
    )
    _write_pptx("from-template", template="acme.potx")

    output_path = _isolate_dirs / "out" / "from-template.pptx"
    output = Presentation(output_path)
    assert output.slide_width == Inches(12)
    with ZipFile(output_path) as package:
        assert b"presentationml.presentation.main+xml" in package.read("[Content_Types].xml")


@pytest.mark.parametrize(
    "writer,template",
    [
        pytest.param(_write_xlsx, "unsafe.xltm", id="xlsx"),
        pytest.param(_write_docx, "unsafe.dotm", id="docx"),
        pytest.param(_write_pptx, "unsafe.potm", id="pptx"),
    ],
)
def test_writers_reject_macro_enabled_templates(writer, template):
    with pytest.raises(ValueError, match="Macro-enabled Office templates"):
        writer("unsafe", template=template)


def test_template_rejects_renamed_embedded_macro(_isolate_dirs):
    template = _isolate_dirs / "in" / "renamed.dotx"
    with ZipFile(template, "w") as package:
        package.writestr("[Content_Types].xml", "<Types />")
        package.writestr("word/vbaProject.bin", b"macro")
    with pytest.raises(ValueError, match="Macro content"):
        _write_docx("unsafe", template="renamed.dotx")


@pytest.mark.parametrize(
    "writer,template",
    [
        pytest.param(_write_xlsx, "outside.xltx", id="xlsx"),
        pytest.param(_write_docx, "outside.dotx", id="docx"),
        pytest.param(_write_pptx, "outside.potx", id="pptx"),
    ],
)
def test_writers_confine_templates(_isolate_dirs, writer, template):
    outside = _isolate_dirs / template
    outside.write_bytes(b"not opened")
    with pytest.raises(FileNotFoundError, match="No Office file found"):
        writer("unsafe", template=str(outside))


def _replace_bytes_in_member(path, member: str, old: bytes, new: bytes) -> None:
    with ZipFile(path) as package:
        payload = package.read(member)
    assert old in payload
    _replace_member(path, member, payload.replace(old, new))
